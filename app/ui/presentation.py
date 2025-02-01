from typing import List, Dict, Tuple
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData, ChartData
import pandas as pd
import logging
import streamlit as st
import time

from app.data.models import Bond, PortfolioConstraints, CreditRating, OptimizationResult, RatingGrade
from app.optimization.engine import OptimizationResult
from app.utils.logging_config import setup_logging


# Configure logger
logger = setup_logging()

# # Add StreamHandler to write logs to streamlit
# class StreamlitHandler(logging.Handler):
#     def emit(self, record):
#         msg = self.format(record)
#         if record.levelno >= logging.ERROR:
#             st.error(msg)
#         elif record.levelno >= logging.WARNING:
#             st.warning(msg)
#         elif record.levelno >= logging.INFO:
#             st.info(msg)
#         else:
#             st.text(msg)

# handler = StreamlitHandler()
# formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
# handler.setFormatter(formatter)
# logger.addHandler(handler)

class PresentationTheme:
    def __init__(
        self,
        title_font: str = "Calibri",
        body_font: str = "Calibri",
        title_size: int = 24,
        subtitle_size: int = 18,
        body_size: int = 12,
        primary_color: Tuple[int, int, int] = (0, 114, 178),  # Blue
        secondary_color: Tuple[int, int, int] = (213, 94, 0),  # Orange
        accent_colors: List[Tuple[int, int, int]] = None
    ):
        self.title_font = title_font
        self.body_font = body_font
        self.title_size = title_size
        self.subtitle_size = subtitle_size
        self.body_size = body_size
        self.primary_color = primary_color
        self.secondary_color = secondary_color
        self.accent_colors = accent_colors or [
            (0, 158, 115),    # Green
            (230, 159, 0),    # Yellow
            (86, 180, 233),   # Light Blue
            (204, 121, 167),  # Pink
            (240, 228, 66),   # Light Yellow
            (0, 114, 178),    # Blue
        ]

def generate_portfolio_presentation(
    result: OptimizationResult,
    universe: List[Bond],
    total_size: float,
    theme: PresentationTheme = None
) -> io.BytesIO:
    """Generate a PowerPoint presentation from optimization results"""
    logger.info("Starting PowerPoint presentation generation")
    theme = theme or PresentationTheme()
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    logger.info("Created presentation with 16:9 aspect ratio")

    # Create portfolio dataframe
    logger.info("Creating portfolio dataframe")
    portfolio_df = pd.DataFrame([
        {
            'ISIN': bond.isin,
            'Issuer': bond.issuer,
            'Country': bond.country,
            'Sector': bond.sector,
            'Rating': bond.credit_rating.display(),
            'Payment Rank': bond.payment_rank,
            'Maturity': bond.maturity_date.strftime('%Y-%m-%d'),
            'Currency': bond.currency,
            'Yield': bond.ytm,
            'Duration': bond.modified_duration,
            'Weight': result.portfolio.get(bond.isin, 0),
            'Market Value': result.portfolio.get(bond.isin, 0) * total_size
        }
        for bond in universe if result.portfolio.get(bond.isin, 0) > 0
    ])
    logger.info(f"Created portfolio dataframe with {len(portfolio_df)} positions")

    # Page 1: Portfolio Characteristics
    logger.info("Generating Page 1: Portfolio Characteristics")
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        title = slide.shapes.title
        title.text = "Portfolio Characteristics"
        
        # Portfolio metrics
        logger.info("Creating characteristics table")
        metrics = {
            'Total Market Value': f"${total_size:,.2f}",
            'Number of Securities': f"{len(portfolio_df)}",
            'Number of Issuers': f"{portfolio_df['Issuer'].nunique()}",
            'Average Duration': f"{result.metrics['duration']:.2f}", 
            'Average Yield': f"{result.metrics['yield']:.2%}", 
            'Average Rating': f"{CreditRating.from_score(float(result.metrics['rating'])).display()}"
        }

        shape = slide.shapes.add_table(rows=len(metrics), cols=2, 
                                     left=Inches(1), top=Inches(2), width=Inches(6), height=Inches(3))
        table = shape.table
        
        for i, (metric, value) in enumerate(metrics.items()):
            table.cell(i, 0).text = metric
            table.cell(i, 1).text = value
        
        # Payment Rank pie chart
        logger.info("Creating payment rank distribution chart")
        payment_rank_data = portfolio_df.groupby('Payment Rank')['Weight'].sum()
        chart_data = CategoryChartData()
        chart_data.categories = payment_rank_data.index
        chart_data.add_series('Payment Rank', [x * 100 for x in payment_rank_data.values])
        x, y = Inches(8), Inches(2)
        cx, cy = Inches(4), Inches(3)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart
        chart.chart_style = 10
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.number_format = '0.0"%"'
        data_labels.font.size = Pt(8)
        data_labels.font.color.rgb = RGBColor(0, 0, 0)
        
        # Ratings bar chart
        logger.info("Creating rating distribution chart")
        ratings_data = portfolio_df.groupby('Rating')['Weight'].sum()
        chart_data = CategoryChartData()
        chart_data.categories = ratings_data.index
        chart_data.add_series('Rating', [x * 100 for x in ratings_data.values])
        x, y = Inches(8), Inches(5)
        cx, cy = Inches(4), Inches(3)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        chart.chart_style = 10
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.number_format = '0.0"%"'
        data_labels.font.size = Pt(8)
        data_labels.font.color.rgb = RGBColor(0, 0, 0)
        chart.value_axis.maximum_scale = 100.0
        chart.value_axis.minimum_scale = 0.0
        chart.value_axis.major_unit = 20.0
        chart.value_axis.minor_unit = 5.0
        chart.value_axis.tick_labels.number_format = '0"%"'
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.has_minor_gridlines = False
        
    except Exception as e:
        logger.error(f"Error creating slide 1: {str(e)}", exc_info=True)
        raise

    # Page 2: Geographic and Sector Analysis
    logger.info("Generating Page 2: Geographic and Sector Analysis")
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Geographic and Sector Analysis"
        
        # Country distribution
        logger.info("Creating country distribution chart")
        country_data = portfolio_df.groupby('Country')['Weight'].sum().sort_values(ascending=True)
        chart_data = CategoryChartData()
        chart_data.categories = country_data.index
        chart_data.add_series('Country', [x * 100 for x in country_data.values])
        x, y = Inches(1), Inches(2)
        cx, cy = Inches(4), Inches(3)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        chart.chart_style = 10
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.number_format = '0.0"%"'
        data_labels.font.size = Pt(8)
        data_labels.font.color.rgb = RGBColor(0, 0, 0)
        chart.value_axis.maximum_scale = 100.0
        chart.value_axis.minimum_scale = 0.0
        chart.value_axis.major_unit = 20.0
        chart.value_axis.minor_unit = 5.0
        chart.value_axis.tick_labels.number_format = '0"%"'
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.has_minor_gridlines = False
        
        # Sector distribution
        logger.info("Creating sector distribution chart")
        sector_data = portfolio_df.groupby('Sector')['Weight'].sum().sort_values(ascending=True)
        chart_data = CategoryChartData()
        chart_data.categories = sector_data.index
        chart_data.add_series('Sector', [x * 100 for x in sector_data.values])
        x, y = Inches(1), Inches(5)
        cx, cy = Inches(4), Inches(3)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        chart.chart_style = 10
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.number_format = '0.0"%"'
        data_labels.font.size = Pt(8)
        data_labels.font.color.rgb = RGBColor(0, 0, 0)
        chart.value_axis.maximum_scale = 100.0
        chart.value_axis.minimum_scale = 0.0
        chart.value_axis.major_unit = 20.0
        chart.value_axis.minor_unit = 5.0
        chart.value_axis.tick_labels.number_format = '0"%"'
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.has_minor_gridlines = False
        
    except Exception as e:
        logger.error(f"Error creating slide 2: {str(e)}", exc_info=True)
        raise

    # Page 3: Top Holdings
    logger.info("Generating Page 3: Top Holdings")
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Top Holdings"
        
        # Top 10 bonds
        logger.info("Creating top 10 bonds table")
        top_bonds = portfolio_df.nlargest(10, 'Weight')[
            ['ISIN', 'Issuer', 'Rating', 'Weight', 'Market Value']
        ].reset_index(drop=True)
        shape = slide.shapes.add_table(rows=len(top_bonds)+1, cols=len(top_bonds.columns), 
                                     left=Inches(1), top=Inches(2), width=Inches(14), height=Inches(2))
        table = shape.table
        # Set headers
        for i, col in enumerate(top_bonds.columns):
            table.cell(0, i).text = col
        # Fill data
        for i in range(len(top_bonds)):
            for j, col in enumerate(top_bonds.columns):
                value = top_bonds.iloc[i][col]
                if col == 'Weight':
                    formatted_value = f"{value:.2%}"
                elif col == 'Market Value':
                    formatted_value = f"${value:,.2f}"
                else:
                    formatted_value = str(value)
                table.cell(i+1, j).text = formatted_value

        # Top 10 issuers
        logger.info("Creating top 10 issuers table")
        top_issuers = portfolio_df.groupby('Issuer')['Weight'].sum().nlargest(10)
        shape = slide.shapes.add_table(rows=len(top_issuers)+1, cols=2, 
                                     left=Inches(1), top=Inches(5), width=Inches(14), height=Inches(2))
        table = shape.table
        # Set headers
        table.cell(0, 0).text = 'Issuer'
        table.cell(0, 1).text = 'Weight'
        # Fill data
        for i, (issuer, weight) in enumerate(top_issuers.items()):
            table.cell(i+1, 0).text = issuer
            table.cell(i+1, 1).text = f"{weight:.2%}"

                    
    except Exception as e:
        logger.error(f"Error creating slide 3: {str(e)}", exc_info=True)
        raise

    # Page 4: Full Portfolio
    logger.info("Generating Page 4: Full Portfolio")
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Full Portfolio"
        
        # Full portfolio table
        logger.info("Creating full portfolio table")
        portfolio_table = portfolio_df[['ISIN', 'Issuer', 'Rating', 'Yield', 'Duration', 'Weight', 'Market Value']].reset_index(drop=True)
        shape = slide.shapes.add_table(rows=len(portfolio_table)+1, cols=len(portfolio_table.columns), 
                                     left=Inches(1), top=Inches(2), width=Inches(14), height=Inches(5))
        table = shape.table
        # Set headers
        for i, col in enumerate(portfolio_table.columns):
            table.cell(0, i).text = col
        # Fill data
        for i in range(len(portfolio_table)):
            for j, col in enumerate(portfolio_table.columns):
                value = portfolio_table.iloc[i][col]
                if col in ['Weight', 'Yield']:
                    formatted_value = f"{value:.2%}"
                elif col == 'Duration':
                    formatted_value = f"{value:.2f}"
                elif col == 'Market Value':
                    formatted_value = f"${value:,.2f}"
                else:
                    formatted_value = str(value)
                table.cell(i+1, j).text = formatted_value
            
    except Exception as e:
        logger.error(f"Error creating slide 4: {str(e)}", exc_info=True)
        raise

    # Save presentation to bytes
    logger.info("Saving presentation to bytes")
    try:
        with io.BytesIO() as pptx_stream:
            prs.save(pptx_stream)
            pptx_stream.seek(0)
            return pptx_stream.getvalue()
    except Exception as e:
        logger.error(f"Error saving presentation: {str(e)}", exc_info=True)
        raise
